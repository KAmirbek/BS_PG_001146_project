import dash
from dash import html, dcc, dash_table, callback
import base64
from dash.dependencies import Input, Output, State
from datetime import datetime, timedelta
import io
import pandas as pd
import sqlite3
import plotly.express as px
import sys
from pptx import Presentation
from pptx.util import Inches
import tempfile
import os
import plotly.io as pio
from dash.exceptions import PreventUpdate
import plotly.graph_objects as go
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from flask import send_file
from pathlib import Path
import numpy as np
from scipy.stats import linregress
from sklearn.linear_model import LinearRegression
from datetime import datetime, timedelta


# Создание экземпляра Dash
app = dash.Dash(__name__)

# Подключение к базе данных
conn = sqlite3.connect('bank_analytics.db')

# Получение списка всех таблиц в базе данных
tables_query = "SELECT name FROM sqlite_master WHERE type='table';"
tables = pd.read_sql_query(tables_query, conn)

# Закрытие соединения с базой данных
conn.close()

# Создание списка для выпадающего списка
table_options = [{'label': table, 'value': table} for table in tables['name']]

# Создание экземпляра Dash
app.layout = html.Div([
    html.H1("Аналитика Банковского сектора", style={'color': 'white', 'text-align': 'center'}),
    # Добавление Tabs
    dcc.Tabs([
        # Вкладка с таблицами
        dcc.Tab(label='Сведения в Таблицах', children=[
            dcc.Dropdown(
                id='table-dropdown',
                options=table_options,
                value=tables['name'][0],  # Начальное значение
            ),
            dash_table.DataTable(
            id='selected-table',
            style_table={'overflowX': 'auto', 'overflowY': 'auto'},
            editable=True,
            filter_action="native",
            sort_action="native",
            style_cell_conditional=[
                {
                    'if': {'column_id': 'any'},  # Применить к любому столбцу
                    'backgroundColor': 'white',  # Задайте цвет фона для ячеек
                },
            ],
        ),
        ]),
        # Вкладка с диаграммами
        dcc.Tab(label='Сведения в Диаграмах по Банковскому сектору', children=[
            dcc.Dropdown(
                id='chart-dropdown',
                options=table_options,
                value=tables['name'][0],  # Начальное значение
            ),
            dcc.Dropdown(
                id='period-dropdown',
                options=[],  # Будет заполнено динамически
                value='',  # Начальный период
            ),
            # Вставляем контейнер для диаграмм
            html.Div(id='charts-container'),
            # В вашем layout внутри вкладки 'Сведения в Диаграмах', добавьте кнопку:
            html.Button("Скачать все диаграммы в PPT", id="btn-download-ppt"),
            dcc.Download(id="download-ppt")
        ]),
        # Вкладка 'Сведения в Диаграмах (Ограниченные Банки)'
        dcc.Tab(label='Сведения в Диаграмах для Государственных Банков', children=[
            dcc.Dropdown(
                id='chart-dropdown-limited-banks',
                options=[{'label': 'Banks', 'value': 'Banks'}],
                value='Banks',  # Фиксированное значение для выбора
                disabled=True,  # Делаем Dropdown неактивным, так как выбор ограничен одним значением
            ),
            dcc.Dropdown(
                id='period-dropdown-limited-banks',
                options=[],  # Будет заполнено динамически
                value='',  # Начальное значение
            ),
            html.Div(id='charts-container-limited-banks'),
        ]),    
        # Вкладка с комбинированными диаграммами для таблицы "Banks"
        dcc.Tab(label='Диаграммы с линией тренда', children=[
            dcc.Dropdown(
                id='combined-banks-period-dropdown',
                options=[],  # Будет заполнено динамически
                multi=True,  # Разрешить выбор нескольких периодов времени
                value=[],  # Начальные периоды
                placeholder="Выберите банк(и)"
            ),
            dcc.Dropdown(
                id='combined-banks-dropdown',
                options=[],  # Будет заполнено динамически
                multi=True,  # Разрешить выбор нескольких банков
                value=[],  # Начальные банки
                placeholder="Выберите период(ы) "
            ),
            # Вставляем контейнер для комбинированных диаграмм
            html.Div(id='combined-banks-charts-container'),
        ]),
        # Код для пятой вкладки,
        dcc.Tab(label='Прогноз на основе линейной регрессии', children=[
            dcc.Dropdown(
                id='forecast-banks-dropdown',  # Обновленный идентификатор
                options=[],  # Будет заполнено динамически
                multi=True,  # Разрешить выбор нескольких банков
                value=[],  # Начальные банки
                placeholder="Выберите банк(и)"
            ),
            # Вставляем контейнер для комбинированных диаграмм
            html.Div(id='forecast-charts-container'),  # Обновленный идентификатор
        ]),
    ]),
], style={'background-image': 'url("/assets/579.jpg")', 'background-size': 'cover', 'height': '100vh', 'overflow': 'auto'})
# Определение обратного вызова для обновления DataTable на основе выбранной вкладки
@app.callback(
    Output('selected-table', 'columns'),
    Output('selected-table', 'data'),
    Input('table-dropdown', 'value')
)
def update_table(selected_table):
    # Подключение к базе данных
    conn = sqlite3.connect('bank_analytics.db')
    try:
        # Чтение данных из выбранной таблицы в DataFrame
        query = f'SELECT * FROM {selected_table};'
        df_selected_table = pd.read_sql_query(query, conn)
        
        # Обработка бинарных данных в DataFrame
        for col in df_selected_table.select_dtypes(include=[object]).columns:
            try:
                df_selected_table[col] = df_selected_table[col].apply(lambda x: x.decode('utf-8') if isinstance(x, bytes) else x)
            except (AttributeError, UnicodeDecodeError) as e:
                print(f"Error decoding column {col}: {e}")       
        
        # Подготовка данных для DataTable
        columns = [{'name': col, 'id': col, 'editable': True} for col in df_selected_table.columns]
        data = df_selected_table.to_dict('records')
        return columns, data
    except Exception as e:
        # Handle any potential exceptions, e.g., if the table does not exist
        print(f"Error: {e}")
        return [], []
# Определение обратного вызова для обновления диаграмм на основе выбранной вкладки
@app.callback(
    Output('charts-container', 'children'),
    Input('chart-dropdown', 'value'),
    Input('period-dropdown', 'value')
)
def update_charts(selected_table, selected_period):
    # Подключение к базе данных
    conn = sqlite3.connect('bank_analytics.db')
    try:
        # Чтение данных из выбранной таблицы в DataFrame
        query = f'SELECT * FROM {selected_table};'
        df_selected_table = pd.read_sql_query(query, conn)
        
        # Исключаем столбцы "Bank_Id" и "Наименование Банка"
        df_selected_table = df_selected_table.drop(columns=["Bank_Id"])
        # Применяем фильтр по выбранному периоду
        # Предполагается, что у вас есть столбец "За период", который можно использовать для фильтрации
        df_selected_table = df_selected_table[df_selected_table['За период'] == selected_period]
        
        # Создаем словарь для назначения цветов
        colors_map = {bank: ('#FFA500' if bank == 'NBU' else '#636EFA') for bank in df_selected_table['Наименование Банка'].unique()}
        # Создаем отдельные столбчатые диаграммы
        charts = []
        
        for col in df_selected_table.columns:
            if col == 'За период':
                continue  # Пропустить столбец "За период" для диаграммы
            if col == 'Наименование Банка':
                continue  # Пропустить столбец "Наименование Банка" для диаграммы
            # Сортируем DataFrame по текущему столбцу в убывающем порядке и используем его для построения диаграммы
            sorted_df = df_selected_table.sort_values(by=col, ascending=False)      
            
            # Создаем диаграмму
            fig = px.bar(
                sorted_df,
                x='Наименование Банка', 
                y=col, 
                title=f'Bar Chart for {col} - {selected_period}',
                labels={'x': 'Банк', 'y':col},
                text=col,
                height=500,
                color='Наименование Банка',
                color_discrete_map=colors_map  # Используем наш словарь для назначения цветов
            )
            
            fig.update_traces(texttemplate='%{text}', textposition='outside')
            # Убираем легенду, поскольку цвета заданы явно
            fig.update_layout(showlegend=False)
            charts.append(dcc.Graph(figure=fig))
        return charts
    except Exception as e:
        # Handle any potential exceptions, e.g., if the table does not exist
        print(f"Error: {e}")
        return []
# Определение обратного вызова для обновления списка периодов в зависимости от выбранной таблицы
@app.callback(
    Output('period-dropdown', 'options'),
    Input('table-dropdown', 'value')
)
def update_period_options(selected_table):
    # Подключение к базе данных
    conn = sqlite3.connect('bank_analytics.db')
    try:
        # Чтение данных из выбранной таблицы в DataFrame
        query = f'SELECT * FROM {selected_table};'
        df_selected_table = pd.read_sql_query(query, conn)
        # Получение уникальных значений из столбца "За период"
        period_options = [
            {"label": period, "value": period} for period in df_selected_table['За период'].unique()
        ]
        return period_options
    except Exception as e:
        # Handle any potential exceptions, e.g., if the table does not exist
        print(f"Error: {e}")
        return []

@app.callback(
    Output("download-ppt", "data"),
    Input("btn-download-ppt", "n_clicks"),
    State('chart-dropdown', 'value'),  # Добавляем State для выбранной таблицы
    State('period-dropdown', 'value'),  # Добавляем State для выбранного периода
    prevent_initial_call=True
)
def download_all_charts_as_ppt(n_clicks, selected_table, selected_period):
    if n_clicks is None:
        raise PreventUpdate
    
    # Подключение к базе данных и получение данных
    conn = sqlite3.connect('bank_analytics.db')
    query = f"SELECT * FROM {selected_table};"
    df_selected_table = pd.read_sql_query(query, conn)
    # Исключаем столбцы "Bank_Id" и "Наименование Банка"
    df_selected_table = df_selected_table.drop(columns=["Bank_Id"])
    # Применяем фильтр по выбранному периоду
    # Предполагается, что у вас есть столбец "За период", который можно использовать для фильтрации
    df_selected_table = df_selected_table[df_selected_table['За период'] == selected_period]
    conn.close()
    # Создание объекта презентации
    prs = Presentation()
    prs.slide_width = Inches(20)
    prs.slide_height = Inches(15)
    
    
    # Путь для сохранения изображений (адаптируйте к вашей ОС)
    output_dir = Path("c:/Users/WINDOWS 10/Desktop/Дипломная/15.04.23/Code/Попытка 15/images")
    output_dir.mkdir(parents=True, exist_ok=True)  # Создаем директорию, если она не существует
    
    def generate_and_add_charts_to_presentation(df_selected_table, selected_period, presentation):
        # Создаем словарь для назначения цветов
        colors_map = {bank: ('#FFA500' if bank == 'NBU' else '#636EFA') for bank in df_selected_table['Наименование Банка'].unique()}
        for col in df_selected_table.columns:
            if col == 'За период':
                continue  # Пропустить столбец "За период" для диаграммы
            if col == 'Наименование Банка':
                continue  # Пропустить столбец "Наименование Банка" для диаграммы
            # Сортируем DataFrame по текущему столбцу в убывающем порядке и используем его для построения диаграммы
            sorted_df = df_selected_table.sort_values(by=col, ascending=False)      
            
            # Создаем диаграмму
            fig = px.bar(
                sorted_df,
                x='Наименование Банка', 
                y=col, 
                title=f'Bar Chart for {col} - {selected_period}',
                labels={'x': 'Банк', 'y':col},
                text=col,
                height=500,
                color='Наименование Банка',
                color_discrete_map=colors_map  # Используем наш словарь для назначения цветов
            )
            
            fig.update_traces(texttemplate='%{text}', textposition='outside')
            # Убираем легенду, поскольку цвета заданы явно
            fig.update_layout(showlegend=False)
            
            # Сохраняем график как изображение
            filename = f"{col}.png"
            # Размеры изображения под А3
            img_path = output_dir / filename
            fig.write_image(str(img_path), scale=5)  # увеличение в 5 раз
            
            # Добавляем новый слайд и изображение в презентацию
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.add_picture(str(img_path), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            # Удаляем временное изображение после добавления его в слайд
            os.remove(img_path)
            
    generate_and_add_charts_to_presentation(df_selected_table, selected_period, prs)
    # Сохраняем презентацию во временный файл
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(temp_file.name)
    temp_file.close()  # Важно закрыть файл перед передачей
    # В Dash, вам нужно вернуть словарь с путем к файлу и именем файла
    return dcc.send_file(temp_file.name)

@app.callback(
    Output('period-dropdown-limited-banks', 'options'),
    Input('chart-dropdown-limited-banks', 'value')
)
def update_period_options_limited_banks(selected_table):
    conn = sqlite3.connect('bank_analytics.db')
    query = f"SELECT DISTINCT `За период` FROM {selected_table};"
    df = pd.read_sql_query(query, conn)
    conn.close()
    return [{'label': period, 'value': period} for period in df['За период']]
@app.callback(
    Output('charts-container-limited-banks', 'children'),
    [Input('chart-dropdown-limited-banks', 'value'),
    Input('period-dropdown-limited-banks', 'value')]
)
def update_charts_limited_banks(selected_table, selected_period):
    conn = sqlite3.connect('bank_analytics.db')
    query = f"SELECT * FROM {selected_table} WHERE `За период` = '{selected_period}' AND `Наименование Банка` IN ('NBU', 'SQB', 'ASAKA', 'AGROBANK', 'XALQ');"
    df = pd.read_sql_query(query, conn)
    conn.close()

    charts = []
    for col in df.columns:
        if col not in ['Bank_Id', 'Наименование Банка', 'За период']:
            sorted_df = df.sort_values(by=col, ascending=False)
            # Преобразуем значения в строки с желаемым форматированием
            texts = [f'{x:.1f}' if x < 1 else f'{x}' for x in sorted_df[col]]
            
            fig = px.bar(sorted_df, x='Наименование Банка', y=col, title=f'Bar Chart for {col} - {selected_period}',
                        text=texts,  # Используем наш список форматированных строк
                        color='Наименование Банка',
                        color_discrete_map={'NBU': '#FFA500', 'SQB': '#636EFA', 'ASAKA': '#636EFA', 'AGROBANK': '#636EFA', 'XALQ': '#636EFA'})
            fig.update_traces(texttemplate='%{text}', textposition='outside')
            fig.update_layout(showlegend=False, yaxis_tickformat = '.2f')
            charts.append(dcc.Graph(figure=fig))

    return charts
    
# Определение обратного вызова для обновления списка периодов и банков в зависимости от выбранной таблицы "Banks"
@app.callback(
    [Output('combined-banks-dropdown', 'options'),
    Output('combined-banks-period-dropdown', 'options')],
    Input('combined-banks-period-dropdown', 'value')
)
def update_combined_banks_options(selected_banks):
    # Подключение к базе данных
    conn = sqlite3.connect('bank_analytics.db')
    try:
        # Чтение данных из таблицы "Banks" в DataFrame
        query = 'SELECT * FROM Banks;'
        df_banks = pd.read_sql_query(query, conn)
        
        # Получение уникальных значений из столбцов "За период" и "Наименование Банка"
        banks_options = [{"label": bank, "value": bank} for bank in df_banks['Наименование Банка'].unique()]
        period_options = [{"label": period, "value": period} for period in df_banks['За период'].unique()]
        
        
        return period_options, banks_options
    except Exception as e:
        # Обработка возможных исключений
        print(f"Error: {e}")
        return [], []

# Определение обратного вызова для обновления комбинированных диаграмм "Banks"
@app.callback(
    Output('combined-banks-charts-container', 'children'),
    [Input('combined-banks-dropdown', 'value'),
    Input('combined-banks-period-dropdown', 'value')]
)
def update_combined_banks_charts(selected_periods, selected_banks):
    # Подключение к базе данных
    conn = sqlite3.connect('bank_analytics.db')
    try:
        # Чтение данных из таблицы "Banks" в DataFrame
        query = 'SELECT * FROM Banks;'
        df_banks = pd.read_sql_query(query, conn)
        
        # Применение фильтров по выбранным периодам и банкам
        df_selected_banks = df_banks[df_banks['За период'].isin(selected_periods) &
                                    df_banks['Наименование Банка'].isin(selected_banks)]

        # Создание комбинированных диаграмм "Banks" (Bar chart + Line graph)
        combined_banks_charts = []
        colors_map = {bank: 'blue' if bank == 'NBU' else None for bank in selected_banks}

        for col in df_selected_banks.columns:
            if col == 'За период':
                continue  # Пропустить столбец "За период" для диаграммы
            if col == 'Наименование Банка':
                continue  # Пропустить столбец "Наименование Банка" для диаграммы
            if col == 'Bank_Id':
                continue  # Пропустить столбец "Наименование Банка" для диаграммы
            
            # Сортировка DataFrame по текущему столбцу
            sorted_df = df_selected_banks.sort_values(by=col, ascending=False)

            # Группировка данных для тренда
            trend_data = (df_selected_banks.groupby('За период')[col]
                          .mean()
                          .reset_index()
                          .sort_values('За период'))
            
            bar_chart = px.bar(
                sorted_df,
                x='За период', 
                y=col, 
                title=f'Combined Chart for {col} - {", ".join(selected_periods)}',
                labels={'x': 'Банк', 'y': col},
                text=col,
                height=500,
                color='Наименование Банка',
                color_discrete_map=colors_map  # Используем наш словарь для назначения цветов
            )
            bar_chart.update_layout(showlegend=False)
            bar_chart.update_traces(texttemplate='%{text}', textposition='outside')
            
            # Line graph
            line_graph = px.line(
                trend_data,
                x=trend_data['За период'],
                y=col,
                title=f'Trend Line for {col} - {", ".join(selected_periods)}',
                labels={'x': 'Банк', 'y': f'{col}'},
                line_shape='linear',
                markers=True,
                height=500
            )
            line_graph.update_traces(line=dict(color='red'))
            
            # Добавление линейного графика тренда к bar chart
            for trace in line_graph['data']:
                bar_chart.add_trace(trace)
            
            combined_banks_charts.append(dcc.Graph(figure=bar_chart))
        
        return combined_banks_charts
    except Exception as e:
        # Обработка возможных исключений
        print(f"Error: {e}")
        return []
# Определение обратного вызова для обновления списка периодов и банков в зависимости от выбранной таблицы "Banks"
@app.callback(
    Output('forecast-banks-dropdown', 'options'),
    [Input('forecast-banks-dropdown', 'value')]  # Или другой триггер, если требуется
)
def update_forecast_banks_options(_):    
    conn = sqlite3.connect('bank_analytics.db')
    try:
        df_banks = pd.read_sql_query('SELECT * FROM Banks;', conn)
        banks_options = [{"label": bank, "value": bank} for bank in df_banks['Наименование Банка'].unique()]

        return banks_options
    except Exception as e:
        # Обработка возможных исключений
        print(f"Error: {e}")
        return [], []

# Определение обратного вызова для обновления комбинированных диаграмм "Banks"
@app.callback(
    Output('forecast-charts-container', 'children'),
    [Input('forecast-banks-dropdown', 'value')]
)
def forecast_banks_charts(selected_banks):
    conn = sqlite3.connect('bank_analytics.db')
    forecast_charts = []

    if not selected_banks:
        return []

    try:
        query = 'SELECT * FROM Banks;'
        df_banks = pd.read_sql_query(query, conn)
        
        # Фильтруем данные только по выбранным банкам
        df_banks = df_banks[df_banks['Наименование Банка'].isin(selected_banks)]
        
        # Преобразуем 'За период' в datetime и сортируем
        df_banks['Дата'] = pd.to_datetime(df_banks['За период'], format='%Y-%m-%d')
        df_banks.sort_values('Дата', inplace=True)

        # Выбор последних 12 месяцев данных
        last_12_months_date = df_banks['Дата'].max() - pd.DateOffset(months=11)
        df_banks = df_banks[df_banks['Дата'] >= last_12_months_date]

        for col in df_banks.columns.drop(['За период', 'Наименование Банка', 'Bank_Id', 'Дата']):
            X = np.arange(len(df_banks)).reshape(-1, 1)
            y = df_banks[col].values

            model = LinearRegression()
            model.fit(X, y)

            future_index = np.arange(0, len(X) + 3).reshape(-1, 1)
            future_preds = model.predict(future_index)

            future_preds_rounded = np.round(future_preds, 2)

            full_dates = pd.date_range(start=df_banks['Дата'].min(), periods=len(future_preds), freq='M')

            fig = px.bar(df_banks, x='Дата', y=col, text=col, labels={'x': 'Дата', 'y': col})
            fig.add_scatter(x=full_dates, y=future_preds_rounded, mode='lines+markers+text', name='Прогноз', line=dict(color='red'), text=future_preds_rounded, textposition="top center", showlegend=False)

            fig.update_layout(title=f'Прогноз для {col} с использованием линейной регрессии')
            fig.update_traces(texttemplate='%{text}')

            forecast_charts.append(dcc.Graph(figure=fig))

        return forecast_charts
    except Exception as e:
        print(f"Ошибка: {e}")
        return []
# Запуск сервера
if __name__ == '__main__':
    app.run_server(debug=True)
